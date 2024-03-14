import pptx

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Length

import matplotlib.pyplot as plt

import numpy as np




def normal_pdf(x_vector, mu=0, sigma=1):
    return (1/(sigma*np.sqrt(2*np.pi)))*np.exp(-((x_vector - mu)/sigma)**2/2)


def is_two_column_text(slide):

    text_title = 0

    if  slide.slide_layout.name != "TITLE":
        all_mu = list()
        all_sigma = list()
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    if shape.has_text_frame:
                        print('SLIDE TITLE: %s'%shape.text_frame.text)
                        text_title = 1
                    continue
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_text_frame:
                centroid_x = shape.left + shape.width / 2
                all_mu.append(Length(centroid_x).mm)
                all_sigma.append(Length(shape.width/4).mm)

            

        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape.text_frame]
        fig_shapes = [shape for shape in slide.shapes if (getattr(shape, 'shape_type')==MSO_SHAPE_TYPE.PICTURE)]
        print("Cantidad text shapes: %d"%(len(text_shapes) - text_title))
        print("Cantidad fig shapes: %d"%len(fig_shapes))
        return (all_mu, all_sigma)
    else:
        return False

# file_path = "C:/Users/daedr/Documents/Docencia_UIS/david_2022/recursos_docencia2022ii/sitio/original/3_SimDigital_28.03.2023.pptx"
file_path = "/home/daedro/Documentos/Dev2024/github_repos/derb_site/original/Simulacion/3_SimDigital_28.03.2023.pptx"

prs = Presentation(file_path)
total_slides = len(prs.slides)
print('Total number of slides: %d'%total_slides)


slide_width = pptx.util.Length(prs.slide_width)
slide_width_emus = slide_width.emu
slide_width_mm = slide_width.mm

print("Slide width in mm: %d"%slide_width_mm)

all_output = list()

for slide_number, slide in enumerate(prs.slides, start=1):
    layout = slide.slide_layout
    print("\n---")
    print("Slide %d uses layout: %s"%(slide_number, layout.name))

    output = is_two_column_text(slide)
    print(output)
        
    # if is_two_column_text(slide):
    #     print("Slide %d migh have two-column text layout"%slide_number)
    # else:
    #     print("No two column text layout for Slide %d"%slide_number)

    all_output.append(output)
        


t_vector = np.arange(1, slide_width_mm)


all_result = list()

for slide_number, output in enumerate(all_output, start=1):
    if output:
        salida = map(lambda mu, sigma: normal_pdf(t_vector, mu, sigma), output[0], output[1])
        result = np.sum(list(salida), axis=0)
        all_result.append(result)

        plt.subplot(7, 2, slide_number)
        plt.plot(t_vector, result)
        plt.title('Slide %d'%slide_number)

plt.show()

# Por cada gaussiana. Generar 30 datos. Posteriormente ver si es posible realizar reconstrucción con un ¿error inferior al 90% en el area? Es necesario hacer gmm? revisar...
# Simplemente necesito minimizar error para definir si son menos columnas que el número de shapes!!


# plt.plot(t_vector, normal_pdf(t_vector))
# plt.show()
