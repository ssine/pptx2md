import pptx
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Length

import matplotlib.pyplot as plt
import numpy as np

from operator import attrgetter

from pptx2md.utils_optim import normal_pdf, f_gauss1, f_gauss2, f_gauss3, fit_column_model, compute_pdf_overlap



def ungroup_shapes(shapes):
  res = []
  for shape in shapes:
    try:
      if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        res.extend(ungroup_shapes(shape.shapes))
      else:
        res.append(shape)
    except Exception as e:
      print(f'failed to load shape {shape}, skipped. error: {e}')
  return res

def is_two_column_text(slide):

    
    if  slide.slide_layout.name != "TITLE":
        all_mu = list()
        all_sigma = list()
        for shape in sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left')):
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    if shape.has_text_frame:
                        print('SLIDE TITLE: %s'%shape.text_frame.text)
                        
                    continue
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_text_frame:
                centroid_x = shape.left + shape.width / 2
                all_mu.append(Length(centroid_x).mm)
                all_sigma.append(Length(shape.width/4).mm) # Gaussiana - 2sigma

        return (all_mu, all_sigma)
    else:
        return False

def assign_shapes(slide, params, ncols=2, slide_width_mm=1000):
    
    shapes_dict = {
        "shapes_pre":list(), 
        "shapes_l":list(), 
        "shapes_c":list(), 
        "shapes_r":list()
    }

    shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))

    print("Ncols is %d"%ncols)

    if ncols==1:
        shapes_dict["shapes_pre"]= sorted(shapes, key = lambda x: x.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER, reverse=True)
        return(shapes_dict)
    elif ncols==2:
        param_means = params[0:2]
        param_sds = params[2:]
    elif ncols==3:
        param_means = params[0:3]
        param_sds = params[3:]
    else:
        raise(ValueError, "Error in the number of columns")
    
    x_vector = np.arange(1, slide_width_mm)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                if shape.has_text_frame:
                    print('SLIDE TITLE: %s'%shape.text_frame.text)
                    shapes_dict["shapes_pre"].insert(0, shape)
                else:
                    shapes_dict["shapes_pre"].append(shape)
                continue
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_text_frame:
            centroid_x = shape.left + shape.width / 2
            curr_mu = Length(centroid_x).mm
            curr_sigma = Length(shape.width/4).mm # Gaussian - 2sigma

            area_u_c = np.zeros(ncols)

            for idx, param_mu in enumerate(param_means):
                area_u_c[idx] = compute_pdf_overlap(normal_pdf(x_vector, mu=param_mu, sigma=param_sds[idx]), 
                                                    normal_pdf(x_vector, curr_mu, curr_sigma))
                
            max_score_column = np.argmax(area_u_c)

            if max_score_column==0:
                shapes_dict["shapes_l"].append(shape)
            elif max_score_column==1:
                if ncols == 2:
                    shapes_dict["shapes_r"].append(shape)
                elif ncols == 3:
                    shapes_dict["shapes_c"].append(shape)
                else:
                    raise(ValueError, "Not allowed number of columns")
            elif max_score_column==2:
                shapes_dict["shapes_r"].append(shape)
            else:
                raise(ValueError, "Max number of columns does not correspond to the number of columns")   

    return(shapes_dict)


if __name__ == "__main__":
    # Testing inference of number of columns
    file_path = "test.pptx"

    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    print('Total number of slides: %d'%total_slides)

    slide_width = pptx.util.Length(prs.slide_width)
    slide_width_emus = slide_width.emu
    slide_width_mm = slide_width.mm

    print("Slide width in mm: %d"%slide_width_mm)

    all_output = list()

    for slide_number, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout
        print("\n---")
        print("Slide %d uses layout: %s"%(slide_number, layout.name))

        output = is_two_column_text(slide)
        print(output)
        
        all_output.append(output)
            
    t_vector = np.arange(1, slide_width_mm)
    all_result = list()

    for slide_number, output in enumerate(all_output, start=1):
        if output:
            print("---")
            print("Slide %d"%slide_number)
            salida = map(lambda mu, sigma: normal_pdf(t_vector, mu, sigma), output[0], output[1])
            result = np.mean(list(salida), axis=0)
            all_result.append(result)

            parameters = fit_column_model(t_vector, result)

            dict_shapes = assign_shapes(prs.slides[slide_number-1], parameters, int(len(parameters)/2), slide_width_mm=slide_width_mm)


            plt.subplot(5, 3, slide_number)
            plt.plot(t_vector, result)
            plt.title('Slide %d'%slide_number)

            if len(parameters)==2:
                plt.plot(t_vector, f_gauss1(t_vector, *parameters), linestyle="dashed")
            elif len(parameters)==4:
                plt.plot(t_vector, f_gauss2(t_vector, *parameters), linestyle="dashed")
            elif len(parameters)==6:
                plt.plot(t_vector, f_gauss3(t_vector, *parameters), linestyle="dashed")

            plt.xlabel(np.array2string(parameters))

    plt.show()